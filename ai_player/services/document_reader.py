from __future__ import annotations

import csv
import hashlib
import json
import re
import shutil
import subprocess
import time
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path

from ai_player.core.config import CONFIG_DIR

SUPPORTED_DOCUMENT_EXTENSIONS = {
    ".csv",
    ".json",
    ".md",
    ".pdf",
    ".pptx",
    ".rtf",
    ".text",
    ".txt",
    ".docx",
}


class DocumentReadCancelled(RuntimeError):
    pass


@dataclass(frozen=True)
class DocumentPage:
    number: int
    title: str
    text: str
    start_seconds: int
    duration_seconds: int
    image_path: str = ""


@dataclass(frozen=True)
class DocumentTranscript:
    source_path: Path
    transcript_path: Path
    title: str
    segment_count: int
    pages: list[DocumentPage]


def document_filter() -> str:
    return (
        "Documents (*.pptx *.docx *.pdf *.txt *.text *.md *.rtf *.csv *.json);;"
        "PowerPoint (*.pptx);;Word (*.docx);;PDF (*.pdf);;Text (*.txt *.text *.md *.rtf);;"
        "Data (*.csv *.json);;All files (*.*)"
    )


def is_supported_document_path(path_value: str | Path) -> bool:
    return Path(path_value).suffix.lower() in SUPPORTED_DOCUMENT_EXTENSIONS


def create_document_transcript(
    path_value: str,
    seconds_per_segment: int = 6,
    cancel_callback: Callable[[], bool] | None = None,
) -> DocumentTranscript:
    path = Path(path_value)
    if not path.exists() or not path.is_file():
        raise RuntimeError("Tep tai lieu khong ton tai.")
    _raise_if_cancelled(cancel_callback)
    pages = read_document_pages(path, seconds_per_segment, cancel_callback=cancel_callback)
    _raise_if_cancelled(cancel_callback)
    if not pages:
        raise RuntimeError("Không trích được nội dung từ tài liệu này.")

    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    transcript_path = CONFIG_DIR / "current_document_transcript.srt"
    entry_index = _write_pages_as_srt(pages, transcript_path)

    return DocumentTranscript(
        source_path=path,
        transcript_path=transcript_path,
        title=path.name,
        segment_count=entry_index - 1,
        pages=pages,
    )


def create_text_document_transcript(
    text: str,
    title: str = "Editor",
    seconds_per_segment: int = 6,
) -> DocumentTranscript:
    clean_text = _clean_text(text)
    if not clean_text:
        raise RuntimeError("Editor chưa có nội dung.")

    blocks = _split_blocks([clean_text], max_chars=220)
    if not blocks:
        blocks = [clean_text]

    pages: list[DocumentPage] = []
    cursor = 0
    for index, block in enumerate(blocks, 1):
        duration = _page_duration_seconds(block, seconds_per_segment)
        pages.append(
            DocumentPage(
                number=index,
                title=f"Trang {index}",
                text=block,
                start_seconds=cursor,
                duration_seconds=duration,
            )
        )
        cursor += duration

    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    transcript_path = CONFIG_DIR / "current_editor_transcript.srt"
    entry_index = _write_pages_as_srt(pages, transcript_path, max_chars=140)

    return DocumentTranscript(
        source_path=transcript_path,
        transcript_path=transcript_path,
        title=title,
        segment_count=entry_index - 1,
        pages=pages,
    )


def read_document_pages(
    path: Path,
    seconds_per_segment: int = 6,
    cancel_callback: Callable[[], bool] | None = None,
) -> list[DocumentPage]:
    _raise_if_cancelled(cancel_callback)
    raw_pages = _read_raw_pages(path)
    _raise_if_cancelled(cancel_callback)
    image_paths = _render_original_pages(path, raw_pages, cancel_callback=cancel_callback)
    _raise_if_cancelled(cancel_callback)
    raw_pages = _align_raw_pages_to_rendered_images(path, raw_pages, len(image_paths))
    if len(image_paths) < len(raw_pages):
        output_dir = _document_render_dir(path)
        image_paths.extend(
            _render_placeholder_pages(
                path,
                raw_pages[len(image_paths) :],
                output_dir,
                start_index=len(image_paths) + 1,
            )
        )
    pages: list[DocumentPage] = []
    cursor = 0
    for index, (number, title, text) in enumerate(raw_pages):
        clean_text = _clean_text(text)
        duration = _page_duration_seconds(clean_text, seconds_per_segment)
        pages.append(
            DocumentPage(
                number=number,
                title=title,
                text=clean_text,
                start_seconds=cursor,
                duration_seconds=duration,
                image_path=str(image_paths[index]) if index < len(image_paths) else "",
            )
        )
        cursor += duration
    return pages


def _write_pages_as_srt(pages: list[DocumentPage], transcript_path: Path, max_chars: int = 220) -> int:
    entry_index = 1
    with transcript_path.open("w", encoding="utf-8") as output:
        for page in pages:
            chunks = _speech_chunks(page.text, max_chars=max_chars)
            if not chunks:
                continue
            chunk_seconds = max(2.0, page.duration_seconds / len(chunks))
            for chunk_index, chunk in enumerate(chunks):
                start = page.start_seconds + int(chunk_index * chunk_seconds)
                if chunk_index == len(chunks) - 1:
                    end = page.start_seconds + page.duration_seconds
                else:
                    end = page.start_seconds + int((chunk_index + 1) * chunk_seconds)
                output.write(f"{entry_index}\n")
                output.write(f"{_srt_time(start)} --> {_srt_time(end)}\n")
                output.write(chunk + "\n\n")
                entry_index += 1
    return entry_index


def _document_render_dir(path: Path) -> Path:
    key_source = f"{path.resolve()}:{path.stat().st_mtime_ns}:{path.stat().st_size}"
    key = hashlib.sha1(key_source.encode("utf-8", errors="ignore")).hexdigest()[:12]
    return CONFIG_DIR / "document_pages" / f"{path.stem}-{key}"


def _read_raw_pages(path: Path) -> list[tuple[int, str, str]]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        pages = _read_docx_pages(path)
    elif suffix == ".pptx":
        pages = _read_pptx_pages(path)
    elif suffix == ".pdf":
        pages = _read_pdf_pages(path)
    elif suffix == ".csv":
        pages = _single_text_page(path, "\n".join(_read_csv(path)))
    elif suffix == ".json":
        pages = _single_text_page(path, "\n".join(_read_json(path)))
    elif suffix == ".rtf":
        pages = _single_text_page(path, "\n".join(_read_rtf(path)))
    elif suffix in {".txt", ".text", ".md"}:
        pages = _text_pages(path)
    elif suffix in {".doc", ".ppt"}:
        raise RuntimeError("Định dạng .doc/.ppt cũ chưa được hỗ trợ trực tiếp. Hãy lưu lại thành .docx/.pptx.")
    else:
        pages = _text_pages(path)
    return pages


def _read_docx_pages(path: Path) -> list[tuple[int, str, str]]:
    from docx import Document

    document = Document(str(path))
    blocks: list[str] = []
    blocks.extend(paragraph.text for paragraph in document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    text = "\n".join(blocks)
    page_blocks = _split_blocks([text])
    return [(index, f"Trang {index}", block) for index, block in enumerate(page_blocks or [""], start=1)]


def _read_pptx_pages(path: Path) -> list[tuple[int, str, str]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[tuple[int, str, str]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_blocks.extend(line.strip() for line in shape.text.splitlines())
        text = "\n".join(item for item in slide_blocks if item)
        pages.append((slide_index, f"Slide {slide_index}", text))
    return pages or [(1, "Slide 1", "")]


def _read_pdf_pages(path: Path) -> list[tuple[int, str, str]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[tuple[int, str, str]] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append((page_index, f"Trang {page_index}", text))
    return pages or [(1, "Trang 1", "")]


def _text_pages(path: Path) -> list[tuple[int, str, str]]:
    blocks = _split_blocks(_read_text(path))
    return [(index, f"Trang {index}", block) for index, block in enumerate(blocks or [""], start=1)]


def _single_text_page(path: Path, text: str) -> list[tuple[int, str, str]]:
    return [(1, path.name, text)]


def _read_csv(path: Path) -> list[str]:
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        return [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]


def _read_json(path: Path) -> list[str]:
    data = json.loads(path.read_text(encoding="utf-8-sig", errors="replace"))
    return _json_text_blocks(data)


def _json_text_blocks(value) -> list[str]:
    if isinstance(value, dict):
        blocks: list[str] = []
        for key, item in value.items():
            child = _json_text_blocks(item)
            if child:
                blocks.extend(f"{key}: {text}" for text in child)
        return blocks
    if isinstance(value, list):
        blocks: list[str] = []
        for item in value:
            blocks.extend(_json_text_blocks(item))
        return blocks
    if value is None:
        return []
    return [str(value)]


def _read_rtf(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text)
    text = re.sub(r"[{}]", " ", text)
    return [text]


def _read_text(path: Path) -> list[str]:
    return path.read_text(encoding="utf-8-sig", errors="replace").splitlines()


def _split_blocks(values: list[str], max_chars: int = 360) -> list[str]:
    blocks: list[str] = []
    for value in values:
        text = _clean_text(value)
        if not text:
            continue
        parts = re.split(r"(?<=[.!?。？！])\s+", text)
        current = ""
        for part in parts:
            part = part.strip()
            if not part:
                continue
            if len(current) + len(part) <= max_chars:
                current = f"{current} {part}".strip()
            else:
                if current:
                    blocks.append(current)
                current = part
        if current:
            blocks.append(current)
    return blocks


def _speech_chunks(text: str, max_chars: int = 220) -> list[str]:
    text = _clean_text(text)
    if not text:
        return []
    parts = re.split(r"(?<=[.!?。？！])\s+", text)
    chunks: list[str] = []
    current = ""
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if len(part) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            chunks.extend(_split_long_text(part, max_chars))
            continue
        candidate = f"{current} {part}".strip()
        if len(candidate) <= max_chars:
            current = candidate
            continue
        if current:
            chunks.append(current)
        current = part
    if current:
        chunks.append(current)
    return chunks


def _split_long_text(text: str, max_chars: int) -> list[str]:
    chunks: list[str] = []
    current = ""
    for word in text.split():
        candidate = f"{current} {word}".strip()
        if len(candidate) <= max_chars:
            current = candidate
            continue
        if current:
            chunks.append(current)
        current = word
    if current:
        chunks.append(current)
    return chunks


def _clean_text(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text


def _srt_time(seconds: int) -> str:
    minutes, sec = divmod(max(0, int(seconds)), 60)
    hours, minutes = divmod(minutes, 60)
    return f"{hours:02d}:{minutes:02d}:{sec:02d},000"


def _page_duration_seconds(text: str, seconds_per_segment: int) -> int:
    if not text.strip():
        return 2
    words = len(text.split())
    estimated = max(int(seconds_per_segment), int(words / 2.4))
    return min(45, max(4, estimated))


def _render_original_pages(
    path: Path,
    raw_pages: list[tuple[int, str, str]],
    cancel_callback: Callable[[], bool] | None = None,
) -> list[Path]:
    _raise_if_cancelled(cancel_callback)
    output_dir = _document_render_dir(path)
    if output_dir.exists():
        shutil.rmtree(output_dir, ignore_errors=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    rendered: list[Path] = []
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        rendered = _render_pdf_pages(path, output_dir, cancel_callback=cancel_callback)
    elif suffix in {".docx", ".pptx"}:
        rendered = _render_office_pages(path, output_dir, cancel_callback=cancel_callback)

    _raise_if_cancelled(cancel_callback)
    if rendered:
        return rendered
    return _render_placeholder_pages(path, raw_pages, output_dir)


def _align_raw_pages_to_rendered_images(
    path: Path,
    raw_pages: list[tuple[int, str, str]],
    image_count: int,
) -> list[tuple[int, str, str]]:
    if path.suffix.lower() != ".docx" or image_count <= 0 or image_count == len(raw_pages):
        return raw_pages

    blocks = [text.strip() for _number, _title, text in raw_pages if text.strip()]
    groups = _balanced_text_groups(blocks, image_count)
    return [(index, f"Trang {index}", "\n".join(group).strip()) for index, group in enumerate(groups, start=1)]


def _balanced_text_groups(blocks: list[str], group_count: int) -> list[list[str]]:
    groups: list[list[str]] = [[] for _ in range(group_count)]
    if not blocks:
        return groups

    total_chars = sum(len(block) for block in blocks)
    target_chars = max(1, total_chars // group_count)
    group_index = 0
    current_chars = 0
    for block_index, block in enumerate(blocks):
        remaining_blocks = len(blocks) - block_index
        remaining_groups = group_count - group_index
        if (
            group_index < group_count - 1
            and groups[group_index]
            and current_chars >= target_chars
            and remaining_blocks >= remaining_groups
        ):
            group_index += 1
            current_chars = 0
        groups[group_index].append(block)
        current_chars += len(block)
    return groups


def _render_pdf_pages(
    path: Path,
    output_dir: Path,
    cancel_callback: Callable[[], bool] | None = None,
) -> list[Path]:
    try:
        import fitz

        document = fitz.open(str(path))
        paths: list[Path] = []
        for page_index, page in enumerate(document, start=1):
            _raise_if_cancelled(cancel_callback)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
            output = output_dir / f"page-{page_index:04d}.png"
            pixmap.save(str(output))
            paths.append(output)
        document.close()
        return paths
    except DocumentReadCancelled:
        raise
    except Exception:
        return []


def _render_office_pages(
    path: Path,
    output_dir: Path,
    cancel_callback: Callable[[], bool] | None = None,
) -> list[Path]:
    soffice = _office_renderer()
    if not soffice:
        return []
    pdf_dir = output_dir / "pdf"
    profile_dir = output_dir / "lo-profile"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    try:
        command = [
            soffice,
            "--headless",
            f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(path),
        ]
        _run_cancelable(command, timeout_seconds=120, cancel_callback=cancel_callback)
    except DocumentReadCancelled:
        raise
    except Exception:
        return []
    pdf_path = pdf_dir / f"{path.stem}.pdf"
    if not pdf_path.exists():
        candidates = list(pdf_dir.glob("*.pdf"))
        pdf_path = candidates[0] if candidates else pdf_path
    return (
        _render_pdf_pages(pdf_path, output_dir, cancel_callback=cancel_callback)
        if pdf_path.exists()
        else []
    )


def _office_renderer() -> str:
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    candidates = [
        Path("C:/Program Files/LibreOffice/program/soffice.exe"),
        Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return ""


def _raise_if_cancelled(cancel_callback: Callable[[], bool] | None) -> None:
    if cancel_callback is not None and cancel_callback():
        raise DocumentReadCancelled("Document read cancelled.")


def _run_cancelable(
    command: list[str],
    *,
    timeout_seconds: float,
    cancel_callback: Callable[[], bool] | None = None,
) -> None:
    startupinfo = None
    if hasattr(subprocess, "STARTUPINFO"):
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    process = subprocess.Popen(
        command,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        startupinfo=startupinfo,
    )
    deadline = time.monotonic() + max(0.1, timeout_seconds)
    try:
        while True:
            _raise_if_cancelled(cancel_callback)
            return_code = process.poll()
            if return_code is not None:
                if return_code != 0:
                    raise subprocess.CalledProcessError(return_code, command)
                return
            if time.monotonic() >= deadline:
                raise subprocess.TimeoutExpired(command, timeout_seconds)
            time.sleep(0.1)
    except Exception:
        _terminate_process(process)
        raise


def _terminate_process(process: subprocess.Popen, timeout_seconds: float = 2.0) -> None:
    if process.poll() is not None:
        return
    try:
        process.terminate()
        process.wait(timeout=timeout_seconds)
    except subprocess.TimeoutExpired:
        process.kill()
        process.wait(timeout=timeout_seconds)
    except Exception:
        try:
            process.kill()
        except Exception:
            pass


def _render_placeholder_pages(
    path: Path,
    raw_pages: list[tuple[int, str, str]],
    output_dir: Path,
    start_index: int = 1,
) -> list[Path]:
    from PIL import Image, ImageDraw

    font_title = _font(34, bold=True)
    font_note = _font(24)
    note = _render_missing_note(path)
    paths: list[Path] = []
    for index, (_number, title, _text) in enumerate(raw_pages, start=start_index):
        image = Image.new("RGB", (1280, 720), "#f8fafc")
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 1279, 719), outline="#d5dee8", width=3)
        draw.text((54, 42), title, fill="#0f172a", font=font_title)
        y = 126
        for line in _wrap_text(note, font_note, 1120, draw):
            draw.text((64, y), line, fill="#475569", font=font_note)
            y += 38
        output = output_dir / f"placeholder-page-{index:04d}.png"
        image.save(output)
        paths.append(output)
    return paths


def _render_missing_note(path: Path) -> str:
    if path.suffix.lower() in {".docx", ".pptx"}:
        return (
            "Không render được ảnh gốc của tài liệu Office. Cài LibreOffice và đảm bảo "
            "lệnh soffice/libreoffice nằm trong PATH để hiển thị đúng slide/trang gốc."
        )
    return "Không render được ảnh gốc của trang này."


def _font(size: int, bold: bool = False):
    from PIL import ImageFont

    candidates = [
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _wrap_text(text: str, font, max_width: int, draw) -> list[str]:
    lines: list[str] = []
    current = ""
    for word in text.split():
        candidate = f"{current} {word}".strip()
        width = draw.textbbox((0, 0), candidate, font=font)[2]
        if width <= max_width:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = word
    if current:
        lines.append(current)
    return lines
